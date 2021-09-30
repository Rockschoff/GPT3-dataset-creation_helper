from PyPDF2 import pdf
import requests
from bs4 import BeautifulSoup

def scrape(l = None):
    print(l)
    links = [
        "https://www.google.com"
    ]
    if l:
        links=l

    soups = []
    file_object = open("prompt.txt" , "w")
    for link in links:
        res = requests.get(link)
        s = BeautifulSoup(res.content , "html.parser")
        soups.append(s)

    # print(len(soups))

    # d= [x.prettify() for x in soups[0].contents]

    for soup in soups:
        body = soup.body
        text = []
        for string in body.strings:
            # print(string)
            s = str(string.encode("utf-8"))
            s=s.strip('\n')
            # print(s[2:len(s) -1])
            
            text.append(s[2:len(s) -1])

        

        for t in text :
            if(len(t)!=2):
                file_object.write(t + "\n")

    # body = soups[0].body

    # text = []
    # for string in body.strings:
    #     # print(string)
    #     s = str(string.encode("utf-8"))
    #     s=s.strip('\n')
    #     # print(s[2:len(s) -1])
        
    #     text.append(s[2:len(s) -1])

    # file_object = open("prompt.txt" , "w")

    # for t in text :
    #     if(len(t)!=2):
    #         file_object.write(t + "\n")

    # this = soups[0].find_all("a")

    import glob

    local_ppts = glob.glob("documents/*pptx")
    local_pdfs = glob.glob("documents/*.pdf")

    pdf_texts=[]
    import PyPDF2
    pdfFileObj = open(local_pdfs[0], 'rb') 
        
    # creating a pdf reader object 
    pdfReader = PyPDF2.PdfFileReader(pdfFileObj) 
        
    # printing number of pages in pdf file 
    num_pages = pdfReader.numPages 
        
    # creating a page object 
    pageObj = pdfReader.getPage(0) 
    for i in range(num_pages):
        page=pdfReader.getPage(i)
        text=page.extractText()
        pdf_texts.append(text)

    # print(pdf_texts)
    # print(len(pdf_texts))

    for t in pdf_texts:
        file_object.write(t)
        
    # closing the pdf file object 
    pdfFileObj.close() 

    from pptx import Presentation

    ppt_texts = []
    for file in local_ppts:
        prs = Presentation(file)
        ppt_text= ""
        for slides in prs.slides:
            
            for shape in slides.shapes:
                if hasattr(shape , "text"):
                    # print("has text")
                    ppt_text = ppt_text + shape.text + "\n"
        ppt_texts.append(ppt_text)

    # print(ppt_texts)      

    for t in ppt_text:
        file_object.write(t)

scrape()


# links = [
#     "https://www.realization.com/why-it-works"
# ]

# soups = []

# for link in links:
#     res = requests.get(link)
#     s = BeautifulSoup(res.content , "html.parser")
#     soups.append(s)

# print(len(soups))

# # d= [x.prettify() for x in soups[0].contents]


# body = soups[0].body

# text = []
# for string in body.strings:
#     # print(string)
#     s = str(string.encode("utf-8"))
#     s=s.strip('\n')
#     # print(s[2:len(s) -1])
    
#     text.append(s[2:len(s) -1])

# file_object = open("prompt.txt" , "w")

# for t in text :
#     if(len(t)!=2):
#         file_object.write(t + "\n")

# this = soups[0].find_all("a")

# import glob

# local_ppts = glob.glob("documents/*pptx")
# local_pdfs = glob.glob("documents/*.pdf")

# pdf_texts=[]
# import PyPDF2
# pdfFileObj = open(local_pdfs[0], 'rb') 
    
# # creating a pdf reader object 
# pdfReader = PyPDF2.PdfFileReader(pdfFileObj) 
    
# # printing number of pages in pdf file 
# num_pages = pdfReader.numPages 
    
# # creating a page object 
# pageObj = pdfReader.getPage(0) 
# for i in range(num_pages):
#     page=pdfReader.getPage(i)
#     text=page.extractText()
#     pdf_texts.append(text)

# print(pdf_texts)
# print(len(pdf_texts))

# for t in pdf_texts:
#     file_object.write(t)
    
# # closing the pdf file object 
# pdfFileObj.close() 

# from pptx import Presentation

# ppt_texts = []
# for file in local_ppts:
#     prs = Presentation(file)
#     ppt_text= ""
#     for slides in prs.slides:
        
#         for shape in slides.shapes:
#             if hasattr(shape , "text"):
#                 print("has text")
#                 ppt_text = ppt_text + shape.text + "\n"
#     ppt_texts.append(ppt_text)

# print(ppt_texts)      

# for t in ppt_text:
#     file_object.write(t)


